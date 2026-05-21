"""
PPTX 课件解析器 — 按幻灯片提取文本（含标题+正文+表格），切片为 RAG 切片。
"""
import os
from pptx import Presentation

MAX_CHARS_PER_SLICE = 2000


def parse_pptx_file(filepath: str) -> list[dict]:
    """解析 PPTX 文件，每张幻灯片提取全部文本作为一个或若干切片。"""
    fn = os.path.basename(filepath)
    try:
        prs = Presentation(filepath)
    except Exception:
        return []

    slices = []
    for slide_idx, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

        text = "\n".join(parts).strip()
        if not text:
            continue

        if len(text) <= MAX_CHARS_PER_SLICE:
            slices.append({
                "file": fn,
                "type": "courseware",
                "content": text,
                "metadata": {
                    "title": f"{fn} slide {slide_idx + 1}",
                    "slide": slide_idx + 1,
                    "filepath": filepath,
                },
            })
        else:
            for chunk_idx in range(0, len(text), MAX_CHARS_PER_SLICE):
                chunk = text[chunk_idx:chunk_idx + MAX_CHARS_PER_SLICE]
                if chunk.strip():
                    slices.append({
                        "file": fn,
                        "type": "courseware",
                        "content": chunk,
                        "metadata": {
                            "title": f"{fn} slide {slide_idx + 1}/{chunk_idx // MAX_CHARS_PER_SLICE + 1}",
                            "slide": slide_idx + 1,
                            "filepath": filepath,
                        },
                    })

    return slices